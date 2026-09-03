#!/usr/bin/env python3
"""
exporter.py — 답을 네 가지 문서로 (18장 "문서로 출력하기")

하나의 중간 형식(markup) 에서 네 갈래로 갈라진다:
    markup = {"title": str, "blocks": [("h", text) | ("p", text) | ("fig", caption, png_bytes) | ("table", caption, rows)]}

    export(markup, "html")  표준 라이브러리만
    export(markup, "docx")  python-docx
    export(markup, "pptx")  python-pptx — 제목마다 슬라이드, 그림은 한 장에 한 슬라이드
    export(markup, "pdf")   reportlab (없으면 HTML 을 돌려주고 경고)

처음엔 출력마다 따로 만들었다가 네 군데가 따로 놀아 중간 형식 하나로 모았다. 갈라지는 지점은 되도록 뒤로.
"""
import html
import io
from dataclasses import dataclass, field


@dataclass
class ExportResult:
    fmt: str
    data: bytes
    mime: str
    warnings: list = field(default_factory=list)


MIME = {"html": "text/html", "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation", "pdf": "application/pdf"}


def answer_to_markup(title, answer_md, figures=(), tables=()):
    """답(마크다운) + 그림·표 → 중간 형식. 제목 줄(#)은 h, 나머지 문단은 p."""
    blocks = []
    for para in [p.strip() for p in answer_md.split("\n\n") if p.strip()]:
        if para.startswith("#"):
            blocks.append(("h", para.lstrip("# ").strip()))
        else:
            blocks.append(("p", para))
    for cap, png in figures:
        blocks.append(("fig", cap, png))
    for cap, rows in tables:
        blocks.append(("table", cap, rows))
    return {"title": title, "blocks": blocks}


def export(markup, fmt="html") -> ExportResult:
    fn = {"html": _html, "docx": _docx, "pptx": _pptx, "pdf": _pdf}[fmt]
    return fn(markup)


def export_all(markup):
    return {fmt: export(markup, fmt) for fmt in ("html", "docx", "pptx", "pdf")}


# ---------------------------------------------------------------- html (항상 된다)

def _html(m) -> ExportResult:
    import base64
    parts = [f"<h1>{html.escape(m['title'])}</h1>"]
    for b in m["blocks"]:
        if b[0] == "h":
            parts.append(f"<h2>{html.escape(b[1])}</h2>")
        elif b[0] == "p":
            parts.append(f"<p>{html.escape(b[1]).replace(chr(10), '<br>')}</p>")
        elif b[0] == "fig":
            src = "data:image/png;base64," + base64.b64encode(b[2]).decode() if b[2] else ""
            parts.append(f"<figure><img src='{src}' style='max-width:100%'><figcaption>{html.escape(b[1])}</figcaption></figure>")
        elif b[0] == "table":
            rows = "".join("<tr>" + "".join(f"<td>{html.escape(str(c))}</td>" for c in r) + "</tr>" for r in b[2])
            parts.append(f"<table border='1'>{rows}</table><p><em>{html.escape(b[1])}</em></p>")
    doc = "<!doctype html><meta charset='utf-8'><body style='font-family:sans-serif;max-width:800px;margin:auto'>" + "\n".join(parts)
    return ExportResult("html", doc.encode("utf-8"), MIME["html"])


# ---------------------------------------------------------------- docx

def _docx(m) -> ExportResult:
    w = []
    try:
        from docx import Document
        from docx.shared import Inches
    except ImportError:
        return ExportResult("docx", b"", MIME["docx"], ["python-docx 가 없다: pip install python-docx"])
    doc = Document()
    doc.add_heading(m["title"], 0)
    for b in m["blocks"]:
        if b[0] == "h":
            doc.add_heading(b[1], 1)
        elif b[0] == "p":
            doc.add_paragraph(b[1])
        elif b[0] == "fig":
            if b[2]:
                doc.add_picture(io.BytesIO(b[2]), width=Inches(5.5))
            doc.add_paragraph(b[1]).italic = True
        elif b[0] == "table":
            rows = b[2]
            if rows:
                t = doc.add_table(rows=len(rows), cols=max(len(r) for r in rows))
                t.style = "Table Grid"
                for i, r in enumerate(rows):
                    for j, c in enumerate(r):
                        t.cell(i, j).text = str(c)
            doc.add_paragraph(b[1])
    buf = io.BytesIO(); doc.save(buf)
    return ExportResult("docx", buf.getvalue(), MIME["docx"], w)


# ---------------------------------------------------------------- pptx

def _pptx(m) -> ExportResult:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return ExportResult("pptx", b"", MIME["pptx"], ["python-pptx 가 없다: pip install python-pptx"])
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0]); s.shapes.title.text = m["title"]
    body_slide, tf = None, None

    def new_slide(heading):
        nonlocal body_slide, tf
        body_slide = prs.slides.add_slide(prs.slide_layouts[1])
        body_slide.shapes.title.text = heading
        tf = body_slide.placeholders[1].text_frame; tf.text = ""
    new_slide(m["title"])
    for b in m["blocks"]:
        if b[0] == "h":
            new_slide(b[1])                                     # 제목마다 슬라이드 하나
        elif b[0] == "p":
            p = tf.add_paragraph(); p.text = b[1][:400]; p.font.size = Pt(14)
        elif b[0] == "fig":                                     # 그림은 한 장에 한 슬라이드
            fs = prs.slides.add_slide(prs.slide_layouts[5]); fs.shapes.title.text = b[1][:80]
            if b[2]:
                fs.shapes.add_picture(io.BytesIO(b[2]), Inches(1), Inches(1.5), width=Inches(8))
        elif b[0] == "table":
            rows = b[2]
            if rows:
                ts = prs.slides.add_slide(prs.slide_layouts[5]); ts.shapes.title.text = b[1][:80]
                shape = ts.shapes.add_table(len(rows), max(len(r) for r in rows), Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * len(rows)))
                for i, r in enumerate(rows):
                    for j, c in enumerate(r):
                        shape.table.cell(i, j).text = str(c)
    buf = io.BytesIO(); prs.save(buf)
    return ExportResult("pptx", buf.getvalue(), MIME["pptx"])


# ---------------------------------------------------------------- pdf

def _pdf(m) -> ExportResult:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
    except ImportError:
        r = _html(m)
        return ExportResult("pdf", r.data, MIME["html"], ["reportlab 이 없다 — HTML 을 대신 돌려준다. pip install reportlab"])
    styles = getSampleStyleSheet()
    font = "Helvetica"
    for path in ("/usr/share/fonts/truetype/nanum/NanumGothic.ttf", "C:/Windows/Fonts/malgun.ttf"):   # 한글 — 없으면 깨진다
        try:
            pdfmetrics.registerFont(TTFont("KR", path)); font = "KR"; break
        except Exception:
            continue
    for st in ("Title", "Heading1", "Normal"):
        styles[st].fontName = font
    story = [Paragraph(html.escape(m["title"]), styles["Title"]), Spacer(1, 12)]
    for b in m["blocks"]:
        if b[0] == "h":
            story.append(Paragraph(html.escape(b[1]), styles["Heading1"]))
        elif b[0] == "p":
            story.append(Paragraph(html.escape(b[1]).replace("\n", "<br/>"), styles["Normal"])); story.append(Spacer(1, 6))
        elif b[0] == "fig":
            if b[2]:
                story.append(Image(io.BytesIO(b[2]), width=400, height=260))
            story.append(Paragraph(html.escape(b[1]), styles["Normal"]))
        elif b[0] == "table":
            if b[2]:
                story.append(Table([[str(c) for c in r] for r in b[2]]))
            story.append(Paragraph(html.escape(b[1]), styles["Normal"]))
    buf = io.BytesIO()
    SimpleDocTemplate(buf, pagesize=A4).build(story)
    return ExportResult("pdf", buf.getvalue(), MIME["pdf"], [] if font == "KR" else ["한글 글꼴을 못 찾았다 — 한글이 깨질 수 있다"])
